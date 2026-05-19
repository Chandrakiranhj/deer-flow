#!/usr/bin/env python3
"""Vision Empower funder report — HTML-first PPTX builder.

A from-scratch rebuild. The previous version hand-drew shapes in python-pptx
and looked like it. This version delegates all visual design to the
html-ppt skill (36 themes, real CSS typography, real gradients, real grids):

    report_data.json
        ↓
    slide manifest (vibe-aware, content-gated)
        ↓
    one HTML file per slide, using html-ppt's class system + theme CSS
        ↓
    headless Chrome screenshots each HTML at 1920×1080
        ↓
    PNGs packed into a 16:9 .pptx, one full-bleed image per slide

Funders open the .pptx in PowerPoint as usual; each slide is a high-res image
rendered with web-quality typography (Inter, Playfair Display, etc.) and the
selected vibe's full visual identity. Slides are no longer editable text, but
funder reports aren't edited after handoff — they're printed and emailed.

Usage:
    python build_pptx.py --data data.json --output out.pptx \\
                         --vibe dark-premium --report-type quarterly \\
                         --period-start 2025-10-01 --period-end 2025-12-31
"""
from __future__ import annotations

import argparse
import base64
import html
import json
import os
import shutil
import subprocess
import sys
import tempfile
import traceback
from datetime import datetime
from io import BytesIO
from pathlib import Path
from urllib.request import Request, urlopen

sys.path.insert(0, str(Path(__file__).resolve().parent))
from _bootstrap import ensure  # type: ignore

ensure({"pptx": "python-pptx", "PIL": "Pillow"})

from pptx import Presentation  # noqa: E402
from pptx.util import Emu, Inches  # noqa: E402
from PIL import Image, ImageOps  # type: ignore  # noqa: E402

from vibes import Vibe, get_vibe  # noqa: E402


# ── Paths ─────────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).resolve().parent
# /mnt/skills/custom/vepip-reports/scripts → /mnt/skills/public/html-ppt
HTML_PPT_DIR = SCRIPT_DIR.parent.parent.parent / "public" / "html-ppt"


# ── Format helpers ────────────────────────────────────────────────────────────
def fmt_money(n) -> str:
    if n is None:
        return "—"
    n = float(n)
    if n >= 1e7:    return f"₹{n / 1e7:.1f}Cr"
    if n >= 1e5:    return f"₹{n / 1e5:.1f}L"
    if n >= 1e3:    return f"₹{n / 1e3:.0f}K"
    return f"₹{int(n):,}"


def fmt_date(d) -> str:
    if not d:
        return "—"
    try:
        dt = datetime.fromisoformat(str(d).replace("Z", "+00:00"))
        return dt.strftime("%-d %b %Y") if sys.platform != "win32" else dt.strftime("%#d %b %Y")
    except Exception:
        return str(d)


def project_duration_months(start: str, end: str) -> int:
    try:
        s = datetime.fromisoformat(str(start).replace("Z", "+00:00"))
        e = datetime.fromisoformat(str(end).replace("Z", "+00:00"))
        return max(1, (e.year - s.year) * 12 + (e.month - s.month))
    except Exception:
        return 0


def word_truncate(s, n: int) -> str:
    if not s:
        return ""
    s = str(s)
    if len(s) <= n:
        return s
    budget = n - 1
    cut = s.rfind(" ", 0, budget)
    if cut < max(8, int(budget * 0.55)):
        cut = budget
    return s[:cut].rstrip(" ,.;:—–-") + "…"


def esc(s) -> str:
    return html.escape("" if s is None else str(s))


def fetch_image(url: str) -> bytes | None:
    if not url:
        return None
    if url.startswith("data:"):
        try:
            _, b64 = url.split(",", 1)
            return base64.b64decode(b64)
        except Exception:
            return None
    try:
        req = Request(url, headers={"User-Agent": "vepip-report-bot/1.0"})
        with urlopen(req, timeout=8) as r:
            return r.read()
    except Exception:
        return None


def img_to_data_uri(data: bytes, target_w: int = 1600,
                    target_aspect: float | None = None) -> str:
    """Center-crop an image (if aspect provided) and inline it as a data URI
    so the rendered HTML carries the image without filesystem refs (works
    across the temp dir without path gymnastics)."""
    try:
        im = Image.open(BytesIO(data))
        im = ImageOps.exif_transpose(im)
        if im.mode not in ("RGB", "RGBA"):
            im = im.convert("RGB")
        if target_aspect:
            sw, sh = im.size
            sa = sw / sh
            if abs(sa - target_aspect) > 0.01:
                if sa > target_aspect:
                    nw = int(sh * target_aspect)
                    left = (sw - nw) // 2
                    im = im.crop((left, 0, left + nw, sh))
                else:
                    nh = int(sw / target_aspect)
                    top = int((sh - nh) * 0.4)
                    im = im.crop((0, top, sw, top + nh))
        if im.size[0] > target_w:
            ratio = target_w / im.size[0]
            im = im.resize((target_w, int(im.size[1] * ratio)), Image.LANCZOS)
        buf = BytesIO()
        if im.mode == "RGBA":
            im = im.convert("RGB")
        im.save(buf, format="JPEG", quality=86, optimize=True)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode("ascii")
    except Exception:
        return ""


# ── Chrome detection ──────────────────────────────────────────────────────────
def _find_repo_root_with_node_modules(start: Path) -> Path | None:
    """Walk up from the script directory to find a directory that contains a
    `node_modules/dom-to-pptx` install. Stops at the filesystem root."""
    p = start.resolve()
    for parent in [p, *p.parents]:
        if (parent / "node_modules" / "dom-to-pptx").exists():
            return parent
    return None


def _node_binary() -> str:
    """Locate the `node` binary. Honour VEPIP_NODE, then PATH lookup, then
    common Windows install locations."""
    env = os.environ.get("VEPIP_NODE")
    if env and Path(env).exists():
        return env
    found = shutil.which("node")
    if found:
        return found
    if sys.platform == "win32":
        candidates = [
            r"C:\Program Files\nodejs\node.exe",
            r"C:\Program Files (x86)\nodejs\node.exe",
        ]
        for c in candidates:
            if Path(c).exists():
                return c
    raise RuntimeError(
        "Node.js not found. Install Node 18+ or set VEPIP_NODE to its executable path."
    )


def find_chrome() -> str | None:
    """Locate a Chromium-class headless browser. Tries the obvious binaries
    on each platform, then falls back to PATH."""
    candidates: list[str] = []
    env = os.environ.get("VEPIP_CHROME")
    if env:
        candidates.append(env)
    if sys.platform == "win32":
        candidates.extend([
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
            r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        ])
    elif sys.platform == "darwin":
        candidates.extend([
            "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
            "/Applications/Chromium.app/Contents/MacOS/Chromium",
            "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        ])
    else:
        # Linux / container runtimes (deer-flow on Linux box)
        for name in ("google-chrome", "google-chrome-stable", "chromium",
                     "chromium-browser", "microsoft-edge"):
            p = shutil.which(name)
            if p:
                candidates.append(p)
        candidates.extend([
            "/usr/bin/google-chrome",
            "/usr/bin/chromium",
            "/usr/bin/chromium-browser",
            "/snap/bin/chromium",
        ])
    for c in candidates:
        if c and Path(c).exists():
            return c
    # PATH fallback
    for name in ("chrome", "chromium", "msedge"):
        p = shutil.which(name)
        if p:
            return p
    return None


def render_html_to_png(html_path: Path, png_path: Path,
                       chrome: str, width: int = 1920, height: int = 1080) -> None:
    """Drive headless Chrome to screenshot one html file at exactly the target
    size. `--virtual-time-budget` advances the browser's virtual clock so
    webfont fetches and layout settle before the screenshot is taken."""
    url = html_path.as_uri()
    args = [
        chrome,
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--no-sandbox",
        "--disable-dev-shm-usage",
        f"--window-size={width},{height}",
        f"--screenshot={png_path}",
        "--virtual-time-budget=8000",
        "--default-background-color=00000000",
        url,
    ]
    try:
        result = subprocess.run(args, capture_output=True, timeout=60)
    except subprocess.TimeoutExpired:
        raise RuntimeError(f"Chrome render timed out for {html_path.name}")
    if not png_path.exists() or png_path.stat().st_size < 1000:
        stderr = result.stderr.decode("utf-8", errors="ignore")[:500]
        raise RuntimeError(
            f"Chrome failed to render {html_path.name}: empty PNG. "
            f"stderr tail: {stderr}"
        )


# ── HTML composer: shell + per-slide layouts ──────────────────────────────────
def _abs_uri(p: Path) -> str:
    return p.resolve().as_uri()


def _shared_css(theme_key: str) -> str:
    """Return <link> tags + <style> block used by both single-slide preview HTML
    and the stacked multi-slide HTML consumed by dom-to-pptx."""
    base_css = _abs_uri(HTML_PPT_DIR / "assets" / "base.css")
    fonts_css = _abs_uri(HTML_PPT_DIR / "assets" / "fonts.css")
    theme_css = _abs_uri(HTML_PPT_DIR / "assets" / "themes" / f"{theme_key}.css")
    return f"""<link rel="stylesheet" href="{fonts_css}">
<link rel="stylesheet" href="{base_css}">
<link rel="stylesheet" id="theme-link" href="{theme_css}">"""


_SECTION_RE = None


def _rewrap_with_id(body: str, slide_id: str) -> str:
    """Take a composer's `<section class="slide is-active" style="...">INNER</section>`
    string and rewrite it so the outer section carries our `id="slide-N"` while
    preserving the composer's own style overrides (e.g. cover sets padding:0).

    Composers were designed for single-slide preview; we re-tag them here so
    dom-to-pptx can address each slide by its section selector.
    """
    import re
    global _SECTION_RE
    if _SECTION_RE is None:
        _SECTION_RE = re.compile(
            r'^\s*<section\b([^>]*)>(.*)</section>\s*$',
            re.IGNORECASE | re.DOTALL,
        )
    m = _SECTION_RE.match(body)
    if not m:
        # Defensive: composer returned bare content — wrap it.
        return (f'<section id="{slide_id}" class="slide ve-slide is-active" '
                f'style="padding:88px 112px;">{body}</section>')
    attrs, inner = m.group(1), m.group(2)
    # Replace `class="slide is-active"` with our enriched class list and
    # inject id="slide-N". Keep any inline style the composer added.
    attrs = re.sub(r'class="[^"]*"', 'class="slide ve-slide is-active"', attrs)
    if 'id=' in attrs:
        attrs = re.sub(r'id="[^"]*"', f'id="{slide_id}"', attrs)
    else:
        attrs = f' id="{slide_id}"' + attrs
    return f'<section{attrs}>{inner}</section>'


def deck_html(theme_key: str, slide_bodies: list[str]) -> str:
    """Stack N slide bodies into ONE HTML document, each as a fixed-size
    1920×1080 section. This is what dom-to-pptx ingests: it walks each section's
    DOM, reads getComputedStyle()/getBoundingClientRect() in pixel space, and
    emits native PPTX shapes — yielding fully editable slides.

    Section IDs are slide-1..slide-N (1-based) — the Node export script passes
    matching selectors to exportToPptx().
    """
    sections = "\n".join(
        _rewrap_with_id(body, f"slide-{i + 1}")
        for i, body in enumerate(slide_bodies)
    )
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>VEPIP deck</title>
{_shared_css(theme_key)}
<style>
  html, body {{ margin: 0; padding: 0; background: #ffffff; }}
  body {{ display: flex; flex-direction: column; gap: 0; }}
  /* Each slide is a fixed 1920×1080 island — dom-to-pptx reads its layout
     from getBoundingClientRect, so dimensions need to be exact. */
  .ve-slide {{ position: relative !important; opacity: 1 !important;
              transform: none !important; flex: none;
              width: 1920px !important; height: 1080px !important;
              padding: 88px 112px;
              box-sizing: border-box !important;
              overflow: hidden !important; }}
  /* Cover slide opts out of padding (it positions absolutely from edges). */
  .ve-slide[data-bleed="true"], .ve-slide.ve-cover {{ padding: 0 !important; }}
  /* All the reusable building blocks the slide bodies reference. */
  .ve-eyebrow {{ font-size: 14px; font-weight: 600; letter-spacing: .22em;
                  text-transform: uppercase; color: var(--accent); }}
  .ve-rule {{ height: 3px; width: 84px; background: var(--accent); border-radius: 2px;
              margin: 22px 0 28px; }}
  .ve-meta {{ color: var(--text-2); font-size: 18px; line-height: 1.55; }}
  .ve-numeral {{ font-family: var(--font-display); font-weight: 800;
                  line-height: .92; letter-spacing: -0.04em; color: var(--text-1); }}
  .ve-card {{ background: var(--surface); border: 1px solid var(--border);
              border-radius: var(--radius); padding: 28px 30px;
              box-shadow: var(--shadow); position: relative; overflow: hidden; }}
  .ve-card-top {{ position: absolute; left: 0; top: 0; right: 0; height: 4px;
                   background: var(--accent); }}
  .ve-progress {{ height: 10px; border-radius: 999px; background: var(--surface-2);
                   overflow: hidden; }}
  .ve-progress > span {{ display: block; height: 100%;
                         background: var(--accent); border-radius: 999px; }}
  .ve-quote-mark {{ font-family: var(--font-serif); font-size: 180px; line-height: .6;
                     color: var(--accent); opacity: .35; }}
  table.ve-table {{ width: 100%; border-collapse: collapse; font-size: 17px; }}
  table.ve-table th {{ text-align: left; font-size: 12px; letter-spacing: .14em;
                        text-transform: uppercase; color: var(--text-3);
                        padding: 0 14px 12px; border-bottom: 1px solid var(--border); }}
  table.ve-table td {{ padding: 16px 14px; border-bottom: 1px solid var(--border);
                        color: var(--text-1); vertical-align: middle; }}
  table.ve-table tr:last-child td {{ border-bottom: 0; }}
  .ve-photo-grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 14px; }}
  .ve-photo-grid img {{ width: 100%; aspect-ratio: 16/10; object-fit: cover;
                         border-radius: var(--radius); }}
</style>
</head>
<body>
{sections}
</body></html>
"""


def slide_html(theme_key: str, body_inner: str,
               extra_style: str = "") -> str:
    """Wrap a single slide's inner HTML in a complete document with html-ppt's
    base CSS + the selected theme. Renders in `body.single` mode so the slide
    fills the viewport with no chrome — used for inspection / debug only."""
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>VEPIP slide</title>
{_shared_css(theme_key)}
<style>
  /* Force exact 16:9 viewport so screenshot matches PPTX 13.333×7.5". */
  html, body {{ width: 1920px; height: 1080px; margin: 0; padding: 0; overflow: hidden; }}
  body.single .slide {{ width: 1920px; height: 1080px; padding: 88px 112px; }}
  /* Reusable building blocks not in base.css */
  .ve-eyebrow {{ font-size: 14px; font-weight: 600; letter-spacing: .22em;
                  text-transform: uppercase; color: var(--accent); }}
  .ve-rule {{ height: 3px; width: 84px; background: var(--accent); border-radius: 2px;
              margin: 22px 0 28px; }}
  .ve-rule.thin {{ height: 2px; opacity: .8; }}
  .ve-meta {{ color: var(--text-2); font-size: 18px; line-height: 1.55;
              letter-spacing: .01em; }}
  .ve-numeral {{ font-family: var(--font-display); font-weight: 800;
                  line-height: .92; letter-spacing: -0.04em; color: var(--text-1); }}
  .ve-card {{ background: var(--surface); border: 1px solid var(--border);
              border-radius: var(--radius); padding: 28px 30px;
              box-shadow: var(--shadow); position: relative; overflow: hidden; }}
  .ve-card-top {{ position: absolute; left: 0; top: 0; right: 0; height: 4px;
                   background: var(--accent); }}
  .ve-progress {{ height: 10px; border-radius: 999px; background: var(--surface-2);
                   overflow: hidden; }}
  .ve-progress > span {{ display: block; height: 100%;
                         background: var(--accent); border-radius: 999px; }}
  .ve-tag {{ display: inline-flex; align-items: center; gap: 8px;
             padding: 6px 14px; border-radius: 999px; font-size: 13px;
             font-weight: 500; background: color-mix(in srgb, var(--accent) 12%, transparent);
             color: var(--accent); border: 1px solid color-mix(in srgb, var(--accent) 28%, transparent); }}
  .ve-grow {{ flex: 1; }}
  .ve-chip-row {{ display: flex; flex-wrap: wrap; gap: 12px; }}
  .ve-quote-mark {{ font-family: var(--font-serif); font-size: 180px; line-height: .6;
                     color: var(--accent); opacity: .35; }}
  table.ve-table {{ width: 100%; border-collapse: collapse; font-size: 17px; }}
  table.ve-table th {{ text-align: left; font-size: 12px; letter-spacing: .14em;
                        text-transform: uppercase; color: var(--text-3);
                        padding: 0 14px 12px; border-bottom: 1px solid var(--border); }}
  table.ve-table td {{ padding: 16px 14px; border-bottom: 1px solid var(--border);
                        color: var(--text-1); vertical-align: middle; }}
  table.ve-table tr:last-child td {{ border-bottom: 0; }}
  .ve-cover-grid {{ display: grid; grid-template-columns: 80px 1fr;
                     height: 100%; gap: 0; }}
  .ve-cover-rail {{ background: var(--accent); height: 100%; }}
  .ve-cover-body {{ padding: 72px 112px 64px; display: flex;
                     flex-direction: column; height: 100%; min-height: 0;
                     box-sizing: border-box; }}
  .ve-photo {{ width: 100%; height: 100%; object-fit: cover; border-radius: var(--radius); }}
  .ve-photo-grid {{ display: grid; grid-template-columns: 1fr 1fr; gap: 14px; }}
  .ve-photo-grid img {{ width: 100%; aspect-ratio: 16/10; object-fit: cover;
                         border-radius: var(--radius); }}
{extra_style}
</style>
</head>
<body class="single">
{body_inner}
</body></html>
"""


def _cover_meta(project, report_type, period_start, period_end):
    is_q = report_type == "quarterly"
    period_str = (
        f"{fmt_date(period_start)} – {fmt_date(period_end)}" if is_q
        else f"{fmt_date(project.get('startDate'))} – {fmt_date(project.get('endDate'))}"
    )
    label = "Quarterly Funder Report" if is_q else "Full Project Report"
    return {
        "period_str": period_str,
        "label": label,
        "name": project.get("name", "Untitled Project"),
        "funder": project.get("funderName", ""),
        "grant": fmt_money(project.get("grantAmount")),
        "states": ", ".join(project.get("states") or []),
    }


def _name_size(n: int, scale: tuple[int, int, int, int] = (84, 104, 128, 144)) -> str:
    """Pick a headline size that fits 2 lines for the given name length.
    Vibes can override the four breakpoints if they want a tighter or
    looser scale."""
    if n > 60:
        return f"{scale[0]}px"
    if n > 40:
        return f"{scale[1]}px"
    if n > 24:
        return f"{scale[2]}px"
    return f"{scale[3]}px"


def compose_cover_editorial(project, *, report_type, period_start, period_end,
                            funder_logo_uri):
    """Editorial-serif cover: warm cream, gold rule, serif italic project
    name, magazine-style metadata column on the right. The most restrained
    of the four covers — designed to age well."""
    m = _cover_meta(project, report_type, period_start, period_end)
    size = _name_size(len(m["name"]), scale=(76, 96, 116, 132))
    logo_html = (f'<img src="{funder_logo_uri}" alt="" '
                 f'style="height:60px;width:auto;opacity:.92;margin-top:12px;">'
                 if funder_logo_uri else "")
    return f"""<section class="slide is-active" style="padding:0;position:relative;">
  <!-- Top double-rule eyebrow band -->
  <div style="position:absolute;left:96px;right:96px;top:80px;">
    <div style="display:flex;justify-content:space-between;align-items:center;">
      <span class="ve-eyebrow" style="color:var(--accent);">{esc(m['label']).upper()}</span>
      <span class="ve-eyebrow" style="color:var(--text-3);">Vision Empower · {esc(fmt_date(datetime.now().isoformat())).upper()}</span>
    </div>
    <div style="height:1px;background:var(--accent);margin-top:18px;"></div>
    <div style="height:1px;background:var(--accent);margin-top:4px;opacity:.4;"></div>
  </div>
  <!-- Project name, magazine-feature treatment -->
  <div style="position:absolute;left:96px;right:520px;top:240px;">
    <div style="font-family:var(--font-serif);font-style:italic;color:var(--text-3);font-size:24px;margin-bottom:24px;letter-spacing:.01em;">{esc(m['period_str'])}</div>
    <h1 style="font-family:var(--font-serif);font-style:italic;font-weight:600;font-size:{size};line-height:0.95;letter-spacing:-0.025em;color:var(--text-1);margin:0;">{esc(m['name'])}</h1>
  </div>
  <!-- Right column: metadata stack, magazine spec sheet -->
  <div style="position:absolute;right:96px;top:240px;width:360px;border-left:1px solid var(--border);padding-left:36px;">
    <div style="margin-bottom:32px;">
      <div class="ve-eyebrow" style="color:var(--text-3);margin-bottom:8px;">Funded by</div>
      <div style="font-family:var(--font-serif);font-size:26px;font-weight:600;color:var(--text-1);line-height:1.2;">{esc(m['funder'])}</div>
      {logo_html}
    </div>
    <div style="margin-bottom:32px;">
      <div class="ve-eyebrow" style="color:var(--text-3);margin-bottom:8px;">Grant value</div>
      <div style="font-family:var(--font-serif);font-size:32px;font-weight:700;color:var(--accent);font-variant-numeric:tabular-nums lining-nums;">{esc(m['grant'])}</div>
    </div>
    {f'<div style="margin-bottom:24px;"><div class="ve-eyebrow" style="color:var(--text-3);margin-bottom:8px;">Operating in</div><div style="font-family:var(--font-serif);font-size:20px;font-weight:500;color:var(--text-1);line-height:1.4;">{esc(m["states"])}</div></div>' if m['states'] else ''}
  </div>
  <!-- Bottom rule + brand mark -->
  <div style="position:absolute;left:96px;right:96px;bottom:80px;">
    <div style="height:1px;background:var(--accent);opacity:.4;"></div>
    <div style="height:1px;background:var(--accent);margin-top:4px;"></div>
    <div style="margin-top:18px;display:flex;justify-content:space-between;align-items:center;">
      <span class="ve-eyebrow" style="color:var(--text-3);">Vision Empower Trust · Enabling Inclusive Education</span>
      <span class="ve-eyebrow" style="color:var(--accent);">visionempower.in</span>
    </div>
  </div>
</section>"""


def compose_cover_dark(project, *, report_type, period_start, period_end,
                      funder_logo_uri):
    """Dark-premium cover: deep navy, massive ghost grant numeral behind the
    title as watermark, accent rail. Confident through scale."""
    m = _cover_meta(project, report_type, period_start, period_end)
    size = _name_size(len(m["name"]), scale=(96, 120, 144, 160))
    logo_html = (f'<img src="{funder_logo_uri}" alt="" '
                 f'style="height:56px;width:auto;opacity:.9;">'
                 if funder_logo_uri else "")
    # Ghost watermark of the grant amount — huge, semi-transparent, anchored
    # right. Reads as texture rather than a number on first glance.
    return f"""<section class="slide is-active ve-cover" style="position:relative;">
  <div style="position:absolute;left:0;top:0;bottom:0;width:80px;background:var(--accent);"></div>
  <!-- Ghost numeral, right edge -->
  <div aria-hidden="true" style="position:absolute;right:-40px;top:120px;font-family:var(--font-display);font-weight:800;font-size:560px;line-height:.85;color:var(--accent);opacity:.06;letter-spacing:-.05em;font-variant-numeric:tabular-nums lining-nums;pointer-events:none;">{esc(m['grant'])}</div>
  <!-- Top eyebrows -->
  <div style="position:absolute;left:144px;top:80px;right:112px;display:flex;justify-content:space-between;align-items:center;">
    <span class="ve-eyebrow">{esc(m['label']).upper()}</span>
    <span class="ve-eyebrow" style="color:var(--text-3);">VISION EMPOWER · {esc(fmt_date(datetime.now().isoformat())).upper()}</span>
  </div>
  <!-- Period + title, mid-vertical -->
  <div style="position:absolute;left:144px;top:240px;right:144px;">
    <div style="color:var(--text-2);font-size:22px;font-weight:500;letter-spacing:.02em;margin-bottom:20px;font-variant-numeric:tabular-nums;">{esc(m['period_str'])}</div>
    <h1 style="font-family:var(--font-display);font-weight:800;font-size:{size};line-height:0.95;letter-spacing:-0.045em;color:var(--text-1);margin:0;">{esc(m['name'])}</h1>
  </div>
  <!-- Bottom: funder block -->
  <div style="position:absolute;left:144px;right:112px;bottom:120px;display:flex;align-items:flex-end;justify-content:space-between;gap:24px;">
    <div>
      <div class="ve-eyebrow" style="color:var(--text-3);margin-bottom:12px;">Funded by</div>
      <div style="font-size:32px;font-weight:700;color:var(--text-1);letter-spacing:-.01em;line-height:1.1;">{esc(m['funder'])}</div>
      <div style="margin-top:14px;color:var(--accent);font-weight:600;font-size:20px;letter-spacing:.01em;font-variant-numeric:tabular-nums;">Grant · {esc(m['grant'])}{'  ·  ' + esc(m['states']) if m['states'] else ''}</div>
    </div>
    <div>{logo_html}</div>
  </div>
</section>"""


def compose_cover_magazine(project, *, report_type, period_start, period_end,
                           funder_logo_uri):
    """Magazine-bold cover: brutalist split. 2/3 saffron block + 1/3 black
    panel. Title in tight all-caps on the saffron, funder details on the
    black. Loud, confident, magazine-cover energy."""
    m = _cover_meta(project, report_type, period_start, period_end)
    name_upper = m["name"].upper()
    size = _name_size(len(name_upper), scale=(72, 96, 120, 144))
    logo_html = (f'<img src="{funder_logo_uri}" alt="" '
                 f'style="height:56px;width:auto;filter:invert(1);opacity:.85;margin-top:14px;">'
                 if funder_logo_uri else "")
    return f"""<section class="slide is-active" style="padding:0;position:relative;display:flex;">
  <!-- Saffron block, 2/3 width -->
  <div style="flex:0 0 1280px;background:var(--accent);color:var(--text-1);position:relative;padding:80px 96px;display:flex;flex-direction:column;">
    <div class="ve-eyebrow" style="color:var(--text-1);opacity:.75;">{esc(m['label']).upper()}</div>
    <div style="margin-top:16px;height:6px;width:96px;background:var(--text-1);"></div>
    <div style="margin-top:48px;color:var(--text-1);font-size:22px;font-weight:700;letter-spacing:.02em;font-variant-numeric:tabular-nums;">{esc(m['period_str'])}</div>
    <h1 style="font-family:var(--font-display);font-weight:900;font-size:{size};line-height:0.92;letter-spacing:-0.04em;text-transform:uppercase;color:var(--text-1);margin:24px 0 0;hyphens:auto;-webkit-hyphens:auto;">{esc(name_upper)}</h1>
    <div style="margin-top:auto;color:var(--text-1);opacity:.7;font-size:14px;font-weight:600;letter-spacing:.18em;text-transform:uppercase;">Vision Empower · {esc(fmt_date(datetime.now().isoformat())).upper()}</div>
  </div>
  <!-- Black panel, 1/3 width -->
  <div style="flex:1;background:var(--text-1);color:var(--surface);padding:80px 64px;display:flex;flex-direction:column;justify-content:center;">
    <div style="color:var(--accent);font-size:14px;font-weight:700;letter-spacing:.18em;text-transform:uppercase;margin-bottom:14px;">Funded by</div>
    <div style="font-family:var(--font-display);font-weight:800;font-size:34px;line-height:1.1;letter-spacing:-0.01em;color:var(--surface);">{esc(m['funder'])}</div>
    {logo_html}
    <div style="margin-top:36px;height:3px;width:48px;background:var(--accent);"></div>
    <div style="margin-top:24px;color:var(--accent);font-size:14px;font-weight:700;letter-spacing:.18em;text-transform:uppercase;">Grant value</div>
    <div style="font-family:var(--font-display);font-weight:800;font-size:48px;line-height:1;color:var(--surface);margin-top:6px;font-variant-numeric:tabular-nums lining-nums;">{esc(m['grant'])}</div>
    {f'<div style="margin-top:36px;height:3px;width:48px;background:var(--accent);"></div><div style="margin-top:24px;color:var(--accent);font-size:14px;font-weight:700;letter-spacing:.18em;text-transform:uppercase;">Operating in</div><div style="color:var(--surface);font-size:18px;font-weight:600;margin-top:6px;line-height:1.4;">{esc(m["states"])}</div>' if m['states'] else ''}
  </div>
</section>"""


def compose_cover_corporate(project, *, report_type, period_start, period_end,
                            funder_logo_uri):
    """Ocean-corporate cover: swiss-grid restraint. Hairline rules, big serial
    "01" in margin, project name as a medium-weight headline (not the loudest
    of the four — corporate covers earn trust through restraint)."""
    m = _cover_meta(project, report_type, period_start, period_end)
    size = _name_size(len(m["name"]), scale=(64, 80, 96, 112))
    logo_html = (f'<img src="{funder_logo_uri}" alt="" '
                 f'style="height:54px;width:auto;opacity:.9;">'
                 if funder_logo_uri else "")
    return f"""<section class="slide is-active" style="padding:0;position:relative;background:var(--bg);">
  <!-- Hairline grid lines (decorative, very subtle) -->
  <div aria-hidden="true" style="position:absolute;inset:0;background-image:linear-gradient(to right,var(--border) 1px,transparent 1px);background-size:160px 100%;opacity:.4;pointer-events:none;"></div>
  <!-- Big serial number top-right -->
  <div style="position:absolute;right:96px;top:80px;text-align:right;">
    <div style="font-family:var(--font-display);font-weight:300;font-size:144px;line-height:.9;color:var(--accent);letter-spacing:-0.04em;font-variant-numeric:tabular-nums lining-nums;">01</div>
    <div class="ve-eyebrow" style="color:var(--text-3);margin-top:-6px;">{esc(m['label']).upper()}</div>
  </div>
  <!-- Hairline horizontal at top -->
  <div style="position:absolute;left:96px;right:96px;top:64px;height:1px;background:var(--text-1);opacity:.85;"></div>
  <!-- Project name, restrained -->
  <div style="position:absolute;left:96px;top:300px;right:480px;">
    <div class="ve-eyebrow" style="color:var(--accent);margin-bottom:18px;">Project</div>
    <h1 style="font-family:var(--font-display);font-weight:600;font-size:{size};line-height:1.02;letter-spacing:-0.025em;color:var(--text-1);margin:0;">{esc(m['name'])}</h1>
    <div style="margin-top:24px;color:var(--text-2);font-size:20px;font-weight:400;letter-spacing:.01em;font-variant-numeric:tabular-nums;">{esc(m['period_str'])}</div>
  </div>
  <!-- 4-column metadata grid at bottom -->
  <div style="position:absolute;left:96px;right:96px;bottom:96px;">
    <div style="height:1px;background:var(--text-1);opacity:.85;"></div>
    <div style="display:grid;grid-template-columns:repeat(4,1fr);gap:24px;margin-top:28px;">
      <div>
        <div class="ve-eyebrow" style="color:var(--text-3);">Funded by</div>
        <div style="font-size:20px;font-weight:600;color:var(--text-1);margin-top:8px;line-height:1.3;">{esc(m['funder'])}</div>
      </div>
      <div>
        <div class="ve-eyebrow" style="color:var(--text-3);">Grant</div>
        <div style="font-size:20px;font-weight:700;color:var(--accent);margin-top:8px;font-variant-numeric:tabular-nums lining-nums;">{esc(m['grant'])}</div>
      </div>
      <div>
        <div class="ve-eyebrow" style="color:var(--text-3);">States</div>
        <div style="font-size:16px;font-weight:500;color:var(--text-1);margin-top:8px;line-height:1.4;">{esc(m['states']) or '—'}</div>
      </div>
      <div style="text-align:right;">
        <div class="ve-eyebrow" style="color:var(--text-3);">Published</div>
        <div style="font-size:16px;font-weight:500;color:var(--text-1);margin-top:8px;font-variant-numeric:tabular-nums;">{esc(fmt_date(datetime.now().isoformat()))}</div>
        {logo_html}
      </div>
    </div>
  </div>
</section>"""


def compose_cover(project: dict, vibe: Vibe, *, report_type: str,
                  period_start: str, period_end: str,
                  funder_logo_uri: str = "") -> str:
    """Dispatch to the per-vibe cover composer. Each vibe gets a structurally
    different cover — not just a recolor — so the same project reads as 4
    distinct documents under the 4 vibes."""
    kw = dict(report_type=report_type, period_start=period_start,
              period_end=period_end, funder_logo_uri=funder_logo_uri)
    key = vibe.key
    if key == "magazine-bold":
        return compose_cover_magazine(project, **kw)
    if key == "dark-premium":
        return compose_cover_dark(project, **kw)
    if key == "ocean-corporate":
        return compose_cover_corporate(project, **kw)
    return compose_cover_editorial(project, **kw)


def compose_overview(project: dict, vibe: Vibe, *, summary: str,
                     period_start: str, period_end: str) -> str:
    months = project_duration_months(period_start, period_end)
    deliv_done = project.get("deliverablesDone") or 0
    deliv_total = project.get("deliverablesTotal") or 0
    cards = [
        (fmt_money(project.get("grantAmount")), "TOTAL GRANT"),
        (str(months) if months else "—", "DURATION (MONTHS)"),
    ]
    if deliv_total:
        cards.append((f"{deliv_done}/{deliv_total}", "DELIVERABLES SO FAR"))

    cards_html = "".join(f"""
      <div class="ve-card" style="text-align:left;">
        <div class="ve-card-top"></div>
        <div class="ve-numeral" style="font-size:64px;">{esc(v)}</div>
        <div class="ve-eyebrow" style="color:var(--text-3);margin-top:14px;">{esc(l)}</div>
      </div>""" for v, l in cards)

    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">01 · Project Overview</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">About this project</h2>
  <div class="ve-rule"></div>
  <div class="grid g2" style="gap:64px;align-items:start;margin-top:16px;">
    <div>
      <p class="lede" style="font-size:24px;line-height:1.55;color:var(--text-2);max-width:none;">{esc(word_truncate(summary, 720))}</p>
      <div style="margin-top:36px;color:var(--text-3);font-size:14px;letter-spacing:.05em;">
        Beneficiaries · Visually impaired children &nbsp;·&nbsp; States · {esc(', '.join(project.get('states') or []))}
      </div>
    </div>
    <div style="display:grid;gap:18px;">
      {cards_html}
    </div>
  </div>
</section>"""
    return inner


def compose_impact(project: dict, vibe: Vibe, *, idx: int) -> str:
    activities = project.get("activities") or []
    teachers = sum(a.get("teachersReached") or 0 for a in activities)
    students = sum(a.get("studentsReached") or 0 for a in activities)
    deliv_done = project.get("deliverablesDone") or 0
    deliv_total = project.get("deliverablesTotal") or 0
    deliv_pct = (
        project.get("deliverablesAvgProgress")
        if project.get("deliverablesAvgProgress") is not None
        else (round(deliv_done / deliv_total * 100) if deliv_total else 0)
    )
    approved = project.get("approvedBudget") or 0
    spent = project.get("spentBudget") or 0
    budget_pct = round(spent / approved * 100) if approved else 0

    cards = [
        (f"{teachers:,}" if teachers else "—", "TEACHERS REACHED",
         "across all field activities"),
        (f"{students:,}" if students else "—", "STUDENTS REACHED",
         "through inclusive classroom kits"),
        (f"{deliv_pct}%", "DELIVERABLE PROGRESS",
         f"{deliv_done} of {deliv_total} on track" if deliv_total else "no targets set"),
        (f"{budget_pct}%", "BUDGET UTILISED",
         f"{fmt_money(spent)} of {fmt_money(approved)}" if approved else "—"),
    ]
    # Hero is the most striking — teachers if present, else deliverables.
    primary_idx = 0 if teachers else 2
    primary = cards[primary_idx]
    others = [c for i, c in enumerate(cards) if i != primary_idx]

    others_html = "".join(f"""
      <div class="ve-card" style="display:flex;align-items:center;gap:22px;">
        <div class="ve-numeral" style="font-size:64px;min-width:170px;">{esc(v)}</div>
        <div>
          <div class="ve-eyebrow" style="color:var(--accent);">{esc(l)}</div>
          <div style="color:var(--text-2);font-size:17px;margin-top:8px;line-height:1.4;">{esc(s)}</div>
        </div>
      </div>""" for v, l, s in others)

    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Key Impact</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">Impact at a glance</h2>
  <div class="ve-rule"></div>
  <div class="grid g2" style="gap:40px;align-items:stretch;margin-top:18px;">
    <div class="ve-card" style="padding:64px;display:flex;flex-direction:column;justify-content:center;">
      <div class="ve-card-top"></div>
      <div class="ve-eyebrow" style="color:var(--accent);">{esc(primary[1])}</div>
      <div class="ve-numeral" style="font-size:240px;margin-top:20px;">{esc(primary[0])}</div>
      <div style="height:3px;width:64px;background:var(--accent);margin:24px 0 18px;border-radius:2px;"></div>
      <div style="color:var(--text-2);font-size:22px;line-height:1.45;">{esc(primary[2])}</div>
    </div>
    <div style="display:grid;gap:18px;align-content:start;">
      {others_html}
    </div>
  </div>
</section>"""
    return inner


def compose_deliverables(project: dict, vibe: Vibe, *, idx: int) -> str:
    items = (project.get("deliverables") or [])[:7]
    done = project.get("deliverablesDone") or 0
    total = project.get("deliverablesTotal") or 0
    rows_html = ""
    for d in items:
        title = esc(d.get("title", ""))
        unit = d.get("unit") or ""
        achieved = d.get("achieved") or 0
        target = d.get("target") or 0
        pct = round(achieved / target * 100) if target else (100 if d.get("status") == "completed" else 0)
        pct_capped = min(pct, 100)
        bar_color = "var(--good)" if pct >= 100 else "var(--accent)" if pct >= 40 else "var(--bad)"
        target_str = f"Target · {target}{' ' + unit if unit else ''}" if target else "No target"
        ach_str = f"{achieved}{' ' + unit if unit else ''}" if achieved else "—"
        rows_html += f"""
        <tr>
          <td style="width:38%;">
            <div style="font-weight:600;color:var(--text-1);font-size:19px;">{title}</div>
            <div style="font-size:14px;color:var(--text-3);margin-top:4px;">{esc(target_str)}</div>
          </td>
          <td style="width:50%;">
            <div class="ve-progress" style="height:14px;">
              <span style="width:{pct_capped}%;background:{bar_color};"></span>
            </div>
            <div style="font-size:13px;color:var(--text-3);margin-top:6px;">{pct}% complete</div>
          </td>
          <td style="width:12%;text-align:right;">
            <div class="ve-numeral" style="font-size:32px;color:{'var(--accent)' if pct >= 100 else 'var(--text-1)'};">{esc(ach_str)}</div>
          </td>
        </tr>"""
    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Deliverables</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">Progress against plan</h2>
  <div class="ve-rule"></div>
  <div style="color:var(--text-3);font-size:15px;margin-bottom:24px;letter-spacing:.04em;">
    {done} OF {total} DELIVERABLES COMPLETE
  </div>
  <table class="ve-table">
    <thead><tr>
      <th>Deliverable</th><th>Progress</th><th style="text-align:right;">Achieved</th>
    </tr></thead>
    <tbody>{rows_html}</tbody>
  </table>
</section>"""
    return inner


def compose_activities(project: dict, vibe: Vibe, *, idx: int,
                        gallery_uris: list[str]) -> str:
    activities = (project.get("activities") or [])[:6]
    has_photos = any(gallery_uris)
    items_html = ""
    for a in activities[:5 if has_photos else 6]:
        title = esc(word_truncate(a.get("title", ""), 64))
        date = fmt_date(a.get("activityDate")) if a.get("activityDate") else ""
        place = a.get("state") or a.get("location") or ""
        meta = "  ·  ".join([p for p in [date, place] if p])
        reaches = []
        if a.get("teachersReached"): reaches.append(f"{a['teachersReached']} teachers")
        if a.get("studentsReached"): reaches.append(f"{a['studentsReached']} students")
        if a.get("schoolsReached"):  reaches.append(f"{a['schoolsReached']} schools")
        reach = "  ·  ".join(reaches) if reaches else esc(word_truncate(a.get("notes") or "", 90))
        items_html += f"""
        <div style="display:flex;gap:18px;align-items:flex-start;margin-bottom:24px;">
          <div style="width:14px;height:14px;border-radius:999px;background:var(--accent);margin-top:7px;flex-shrink:0;"></div>
          <div style="flex:1;">
            <div style="font-weight:600;color:var(--text-1);font-size:19px;line-height:1.35;">{title}</div>
            {f'<div style="color:var(--accent);font-size:13px;margin-top:4px;letter-spacing:.04em;">{esc(meta)}</div>' if meta else ''}
            {f'<div style="color:var(--text-2);font-size:16px;margin-top:6px;">{reach}</div>' if reach else ''}
          </div>
        </div>"""

    if has_photos:
        photos = [u for u in gallery_uris if u][:4]
        # 2x2 if 4, 1x2 stacked if fewer
        photos_html = '<div class="ve-photo-grid">' + "".join(
            f'<img src="{u}" alt="">' for u in photos
        ) + '</div>'
        layout = f"""
  <div class="grid g2" style="gap:48px;align-items:start;margin-top:18px;">
    <div>{items_html}</div>
    <div>{photos_html}</div>
  </div>"""
    else:
        layout = f'<div style="margin-top:18px;">{items_html}</div>'

    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Activities</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">Field activities & reach</h2>
  <div class="ve-rule"></div>
  {layout}
</section>"""
    return inner


def compose_stories(project: dict, vibe: Vibe, *, idx: int,
                    testimonials: list[dict]) -> str:
    strong = [t for t in testimonials if (t.get("content") or "").strip()]
    if not strong:
        return ""
    primary = strong[0]
    secondary = strong[1] if len(strong) > 1 else None
    primary_role = primary.get("role") or ""
    primary_attrib = f"{primary.get('author', '')}{' · ' + primary_role if primary_role else ''}"

    secondary_html = ""
    if secondary:
        sec_role = secondary.get("role") or ""
        sec_attrib = f"{secondary.get('author', '')}{', ' + sec_role if sec_role else ''}"
        secondary_html = f"""
    <div style="margin-top:48px;padding-top:36px;border-top:1px solid var(--border);display:flex;gap:24px;align-items:flex-start;">
      <div style="font-family:var(--font-serif);font-size:80px;line-height:.5;color:var(--accent);opacity:.4;">&ldquo;</div>
      <div style="flex:1;">
        <div style="font-size:22px;font-style:italic;color:var(--text-2);line-height:1.5;font-family:var(--font-serif);">{esc(word_truncate(secondary.get('content', ''), 220))}</div>
        <div style="margin-top:14px;color:var(--text-3);font-size:14px;letter-spacing:.04em;">— {esc(sec_attrib)}</div>
      </div>
    </div>"""

    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Stories from the field</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">In their own words</h2>
  <div class="ve-rule"></div>
  <div style="margin-top:32px;">
    <div class="ve-quote-mark">&ldquo;</div>
    <div style="font-family:var(--font-serif);font-size:42px;line-height:1.35;font-style:italic;color:var(--text-1);margin-top:-24px;max-width:1400px;letter-spacing:-0.01em;">{esc(word_truncate(primary.get('content', ''), 360))}</div>
    <div style="margin-top:32px;color:var(--accent);font-weight:600;letter-spacing:.04em;font-size:16px;text-transform:uppercase;">— {esc(primary_attrib)}</div>
    {secondary_html}
  </div>
</section>"""
    return inner


def compose_geographic(project: dict, vibe: Vibe, *, idx: int) -> str:
    states = project.get("states") or []
    activities = project.get("activities") or []
    teachers = sum(a.get("teachersReached") or 0 for a in activities)
    students = sum(a.get("studentsReached") or 0 for a in activities)
    schools  = sum(a.get("schoolsReached") or 0 for a in activities)

    state_chips = "".join(
        f'<div class="ve-card" style="padding:18px 24px;border-left:3px solid var(--accent);">'
        f'<div style="color:var(--text-1);font-weight:600;font-size:20px;">{esc(s)}</div>'
        f'</div>' for s in states
    )
    metrics = [
        (f"{teachers:,}" if teachers else "—", "Teachers"),
        (f"{students:,}" if students else "—", "Students"),
        (f"{schools:,}" if schools else "—", "Schools"),
        (str(len(activities)), "Activities"),
    ]
    metrics_html = "".join(f"""
      <div>
        <div class="ve-numeral" style="font-size:48px;color:var(--accent);">{esc(v)}</div>
        <div class="ve-eyebrow" style="color:var(--text-3);margin-top:6px;">{esc(l)}</div>
      </div>""" for v, l in metrics)

    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Geographic reach</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">Where we work</h2>
  <div class="ve-rule"></div>
  <div class="grid g2" style="gap:64px;align-items:start;margin-top:18px;">
    <div style="text-align:center;display:flex;flex-direction:column;justify-content:center;height:100%;">
      <div class="ve-numeral" style="font-size:260px;color:var(--accent);">{len(states)}</div>
      <div class="ve-eyebrow" style="color:var(--text-3);margin-top:14px;font-size:16px;">STATES COVERED</div>
    </div>
    <div>
      <div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin-bottom:36px;">{state_chips}</div>
      <div style="display:grid;grid-template-columns:repeat(4,1fr);gap:24px;padding-top:28px;border-top:1px solid var(--border);">{metrics_html}</div>
    </div>
  </div>
</section>"""
    return inner


def _svg_donut(pct: int, used_color: str = "var(--accent)",
               rest_color: str = "var(--border-strong)", size: int = 360) -> str:
    """Inline SVG donut. Crisp at any zoom, themable via CSS vars."""
    pct = max(0, min(100, int(pct)))
    r = 130
    cx = cy = size // 2
    stroke = 36
    circ = 2 * 3.141592653589793 * r
    dash = circ * pct / 100
    rest = circ - dash
    return f"""<svg width="{size}" height="{size}" viewBox="0 0 {size} {size}">
  <circle cx="{cx}" cy="{cy}" r="{r}" fill="none" stroke="{rest_color}" stroke-width="{stroke}"/>
  <circle cx="{cx}" cy="{cy}" r="{r}" fill="none" stroke="{used_color}"
          stroke-width="{stroke}" stroke-dasharray="{dash} {rest}"
          stroke-dashoffset="{circ / 4}" stroke-linecap="butt"
          transform="rotate(-90 {cx} {cy})"/>
</svg>"""


def compose_financials(project: dict, vibe: Vibe, *, idx: int) -> str:
    approved = project.get("approvedBudget") or 0
    spent = project.get("spentBudget") or 0
    balance = max(approved - spent, 0)
    util_pct = round(spent / approved * 100) if approved else 0
    used_color = "var(--bad)" if util_pct >= 95 else "var(--accent)"
    donut_html = _svg_donut(util_pct, used_color=used_color)

    metric_cards = [
        (fmt_money(approved), "Approved budget"),
        (fmt_money(spent), "Amount spent"),
        (fmt_money(balance), "Balance remaining"),
    ]
    metrics_html = "".join(f"""
      <div style="padding:24px 0;border-top:1px solid var(--border);">
        <div class="ve-numeral" style="font-size:54px;">{esc(v)}</div>
        <div style="color:var(--text-3);font-size:14px;letter-spacing:.06em;margin-top:6px;text-transform:uppercase;">{esc(l)}</div>
      </div>""" for v, l in metric_cards)

    budgets = (project.get("budgets") or [])[:4]
    bars_html = ""
    if budgets:
        # Find max approved for x-axis normalisation
        max_appr = max((float(b.get("approvedAmount") or 0) for b in budgets), default=1)
        bars_rows = ""
        for b in budgets:
            name = esc(word_truncate(b.get("name", ""), 24))
            ap = float(b.get("approvedAmount") or 0)
            sp = float(b.get("spentAmount") or 0)
            ap_pct = (ap / max_appr * 100) if max_appr else 0
            sp_pct = (sp / max_appr * 100) if max_appr else 0
            sp_color = "var(--bad)" if (ap and sp / ap >= 0.95) else "var(--accent)"
            bars_rows += f"""
            <div style="margin-bottom:20px;">
              <div style="display:flex;justify-content:space-between;align-items:baseline;margin-bottom:8px;">
                <div style="font-weight:600;color:var(--text-1);font-size:17px;">{name}</div>
                <div style="color:var(--text-3);font-size:13px;">{esc(fmt_money(sp))} / {esc(fmt_money(ap))}</div>
              </div>
              <div style="position:relative;height:14px;background:var(--surface-2);border-radius:999px;overflow:hidden;">
                <div style="position:absolute;left:0;top:0;width:{ap_pct}%;height:100%;background:color-mix(in srgb, {sp_color} 18%, transparent);"></div>
                <div style="position:absolute;left:0;top:0;width:{sp_pct}%;height:100%;background:{sp_color};border-radius:999px;"></div>
              </div>
            </div>"""
        bars_html = f"""
      <div style="margin-top:38px;">
        <div class="ve-eyebrow" style="color:var(--text-3);font-size:12px;margin-bottom:18px;">BUDGET BREAKDOWN BY CATEGORY</div>
        {bars_rows}
      </div>"""

    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Financials</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">Budget & utilisation</h2>
  <div class="ve-rule"></div>
  <div class="grid g2" style="gap:64px;align-items:start;margin-top:18px;">
    <div style="position:relative;display:flex;align-items:center;justify-content:center;height:360px;">
      {donut_html}
      <div style="position:absolute;text-align:center;">
        <div class="ve-numeral" style="font-size:88px;">{util_pct}%</div>
        <div class="ve-eyebrow" style="color:var(--text-3);margin-top:6px;">UTILISED</div>
      </div>
    </div>
    <div>{metrics_html}</div>
  </div>
  {bars_html}
</section>"""
    return inner


def compose_way_forward(project: dict, vibe: Vibe, *, idx: int,
                        challenges: str, next_steps: str) -> str:
    chal_html = ""
    if challenges:
        chal_html = f"""
    <div class="ve-card" style="background:var(--surface-2);border-left:4px solid var(--bad);">
      <div class="ve-eyebrow" style="color:var(--bad);font-size:13px;">CHALLENGES</div>
      <h3 style="margin-top:8px;margin-bottom:14px;font-size:26px;letter-spacing:-0.01em;">What got in the way</h3>
      <div style="color:var(--text-2);font-size:18px;line-height:1.55;">{esc(word_truncate(challenges, 520))}</div>
    </div>"""
    next_html = ""
    if next_steps:
        next_html = f"""
    <div class="ve-card" style="background:var(--surface-2);border-left:4px solid var(--good);">
      <div class="ve-eyebrow" style="color:var(--good);font-size:13px;">NEXT STEPS</div>
      <h3 style="margin-top:8px;margin-bottom:14px;font-size:26px;letter-spacing:-0.01em;">Where we go from here</h3>
      <div style="color:var(--text-2);font-size:18px;line-height:1.55;">{esc(word_truncate(next_steps, 520))}</div>
    </div>"""
    cols = "1fr 1fr" if (chal_html and next_html) else "1fr"
    inner = f"""<section class="slide is-active">
  <span class="ve-eyebrow">0{idx} · Way forward</span>
  <h2 class="h2" style="margin-top:14px;font-size:64px;">Honest reflection</h2>
  <div class="ve-rule"></div>
  <div style="display:grid;grid-template-columns:{cols};gap:36px;margin-top:18px;">
    {next_html}
    {chal_html}
  </div>
</section>"""
    return inner


def compose_closing(project: dict, vibe: Vibe, *, idx: int) -> str:
    name = esc(project.get("name", ""))
    funder = esc(project.get("funderName", ""))
    inner = f"""<section class="slide is-active" style="text-align:center;justify-content:center;">
  <div style="max-width:1300px;margin:0 auto;">
    <span class="ve-eyebrow" style="display:block;margin-bottom:32px;">A note of gratitude</span>
    <h1 class="h1" style="font-size:200px;font-family:var(--font-serif);font-style:italic;line-height:.95;letter-spacing:-0.04em;font-weight:600;">Thank you.</h1>
    <div style="height:3px;width:120px;background:var(--accent);margin:48px auto;border-radius:2px;"></div>
    <div style="font-size:24px;color:var(--text-2);font-weight:500;">{name}</div>
    <div style="margin-top:10px;color:var(--text-3);font-size:18px;">Funded by {funder}</div>
    <div style="margin-top:64px;color:var(--accent);font-weight:600;letter-spacing:.18em;font-size:14px;text-transform:uppercase;">
      Vision Empower Trust · visionempower.in
    </div>
  </div>
</section>"""
    return inner


# ── PPTX packer ───────────────────────────────────────────────────────────────
def pack_pngs_to_pptx(png_paths: list[Path], out_path: Path,
                      project_name: str) -> None:
    """Pack N PNGs into a 13.333×7.5 inch .pptx, one full-bleed image per
    slide. python-pptx is great at exactly this — full-bleed image, no shape
    drawing required."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = f"{project_name} — Funder Report"
    prs.core_properties.author = "Vision Empower"
    blank_layout = prs.slide_layouts[6]  # blank layout
    for png in png_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png), Emu(0), Emu(0),
            width=prs.slide_width, height=prs.slide_height,
        )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


# ── Slide manifest (vibe-aware, content-gated) ────────────────────────────────
def build_manifest(project: dict, vibe: Vibe, draft: str) -> list[str]:
    """Decide which slide types this project earns. Vibe knobs tune the
    threshold (prefers_prose lowers the bar for overview/way-forward;
    prefers_quotes lowers it for stories; etc.)."""
    narrative = project.get("narrative") if isinstance(project.get("narrative"), dict) else {}
    summary_text = ((narrative.get("overview") or "").strip() if narrative else "") or project.get("summary") or ""
    deliverables = project.get("deliverables") or []
    activities = project.get("activities") or []
    testimonials = list(project.get("testimonials") or [])
    if not testimonials:
        for a in activities:
            if a.get("testimonial"):
                testimonials.append({
                    "content": a["testimonial"],
                    "author": a.get("testimonialBy") or a.get("title", ""),
                    "role": a.get("state") or "",
                })
    states = project.get("states") or []
    approved = project.get("approvedBudget") or 0
    spent = project.get("spentBudget") or 0
    teachers = sum(a.get("teachersReached") or 0 for a in activities)
    students = sum(a.get("studentsReached") or 0 for a in activities)
    has_narrative_forward = (
        (narrative.get("challenges") or "").strip()
        or (narrative.get("way_forward") or "").strip()
    ) if narrative else False

    quote_threshold = 40 if vibe.layout.prefers_quotes else 100
    overview_min = 80 if vibe.layout.prefers_prose else 180
    way_forward_min_draft = 300 if vibe.layout.prefers_prose else 600

    gates = [
        ("cover", True),
        ("overview", bool(summary_text and len(summary_text) >= overview_min) or bool(draft and len(draft) >= max(200, overview_min * 2))),
        ("impact", bool(deliverables) or bool(approved) or teachers > 0 or students > 0),
        ("deliverables", len(deliverables) >= 1),
        ("activities", len(activities) >= 1),
        ("stories", any((t.get("content") or "").strip() and len((t.get("content") or "").strip()) >= quote_threshold for t in testimonials)),
        ("geographic", len(states) >= 2),
        ("financials", approved > 0 or (vibe.layout.prefers_charts and (spent > 0 or bool(project.get("budgets"))))),
        ("way_forward", has_narrative_forward or bool(draft and len(draft) >= way_forward_min_draft)),
        ("closing", True),
    ]
    manifest = [name for name, ok in gates if ok]

    # Apply slide_count_cap by trimming least-essential slides first.
    drop_order = ["geographic", "overview", "way_forward", "stories",
                  "impact", "financials", "deliverables", "activities"]
    cap = max(4, vibe.layout.slide_count_cap)
    for name in drop_order:
        if len(manifest) <= cap:
            break
        if name in manifest:
            manifest.remove(name)
    return manifest


# ── Main ──────────────────────────────────────────────────────────────────────
def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--report-type", default="quarterly", choices=["quarterly", "full"])
    parser.add_argument("--period-start", default="")
    parser.add_argument("--period-end", default="")
    parser.add_argument("--draft", default="")
    parser.add_argument("--vibe", default="editorial-serif",
                        help="editorial-serif | dark-premium | magazine-bold | ocean-corporate")
    parser.add_argument("--keep-html", action="store_true",
                        help="Keep intermediate HTML/PNG files for debugging")
    args = parser.parse_args(argv)

    chrome = find_chrome()
    if not chrome:
        raise RuntimeError(
            "No headless browser found. Install Google Chrome, Chromium, or "
            "Microsoft Edge, or set VEPIP_CHROME to its executable path."
        )

    if not HTML_PPT_DIR.exists():
        raise RuntimeError(
            f"html-ppt skill not found at {HTML_PPT_DIR}. "
            "This builder depends on the html-ppt skill being installed alongside vepip-reports."
        )

    data = json.loads(Path(args.data).read_text(encoding="utf-8"))
    project = data.get("project") or data

    # Normalise shape variants from upstream
    if "budgets" not in project and project.get("budgetCategories"):
        project["budgets"] = project["budgetCategories"]
    if "activities" not in project and project.get("recentActivities"):
        project["activities"] = project["recentActivities"]
    if "deliverablesTotal" not in project or "deliverablesDone" not in project or "deliverablesAvgProgress" not in project:
        delivs = project.get("deliverables") or []
        project.setdefault("deliverablesTotal", len(delivs))
        project.setdefault(
            "deliverablesDone",
            sum(1 for d in delivs if (d.get("achieved") or 0) >= (d.get("target") or 0) > 0),
        )
        if delivs:
            ratios = [
                min(1.0, (d.get("achieved") or 0) / (d.get("target") or 0))
                for d in delivs if (d.get("target") or 0) > 0
            ]
            avg = round(sum(ratios) / len(ratios) * 100) if ratios else 0
        else:
            avg = 0
        project.setdefault("deliverablesAvgProgress", avg)
    if "approvedBudget" not in project or "spentBudget" not in project:
        cats = project.get("budgetCategories") or project.get("budgets") or []
        project.setdefault("approvedBudget",
                          sum((c.get("approvedAmount") or c.get("approved") or 0) for c in cats))
        project.setdefault("spentBudget",
                          sum((c.get("spentAmount") or c.get("spent") or 0) for c in cats))

    period_start = args.period_start or data.get("periodStart") or project.get("startDate") or ""
    period_end   = args.period_end   or data.get("periodEnd")   or project.get("endDate")   or ""
    draft        = args.draft or data.get("draft") or ""

    vibe = get_vibe(args.vibe)
    theme = vibe.html_theme or "editorial-serif"

    # Fetch + crop assets so they're embedded in the HTML as data URIs.
    funder_logo_uri = ""
    if project.get("funderLogoUrl"):
        img = fetch_image(project["funderLogoUrl"])
        if img:
            funder_logo_uri = img_to_data_uri(img, target_w=240)
    gallery_uris: list[str] = []
    for g in (project.get("gallery") or [])[:4]:
        img = fetch_image((g or {}).get("url", ""))
        if img:
            gallery_uris.append(img_to_data_uri(img, target_w=1600, target_aspect=16/10))

    # Build narrative text fields for slides
    narrative = project.get("narrative") if isinstance(project.get("narrative"), dict) else {}
    summary = ""
    if isinstance(narrative, dict):
        summary = (narrative.get("overview") or "").strip()
    if not summary:
        summary = project.get("summary") or ""
    if not summary and draft:
        lines = [l.strip() for l in draft.split("\n") if len(l.strip()) > 60]
        summary = " ".join(lines[:4])[:600]

    testimonials = list(project.get("testimonials") or [])
    if not testimonials:
        for a in (project.get("activities") or []):
            if a.get("testimonial"):
                testimonials.append({
                    "content": a["testimonial"],
                    "author": a.get("testimonialBy") or a.get("title", ""),
                    "role": a.get("state") or "",
                })

    challenges = (narrative.get("challenges") or "").strip() if narrative else ""
    next_steps = (narrative.get("way_forward") or "").strip() if narrative else ""
    if (not challenges or not next_steps) and draft:
        import re as _re
        def extract(text, start_re, stop_re):
            parts = _re.split(start_re, text, flags=_re.IGNORECASE, maxsplit=1)
            if len(parts) < 2: return ""
            content = parts[1]
            m = _re.search(stop_re, content, flags=_re.IGNORECASE)
            if m: content = content[: m.start()]
            return _re.sub(r"\s+", " ", _re.sub(r"^[\s:–\-]+", "", content)).strip()[:600]
        if not challenges:
            challenges = extract(draft, r"challenges?|barriers?|constraints?",
                                 r"next\s+steps?|way\s+forward|financial")
        if not next_steps:
            next_steps = extract(draft, r"next\s+steps?|way\s+forward|upcoming\s+activities",
                                 r"evidence|conclusion|thank|challenges?")

    manifest = build_manifest(project, vibe, draft)

    # Renderer scratchpad
    keep = args.keep_html
    tmp = Path(tempfile.mkdtemp(prefix="vepip_pptx_"))
    if keep:
        print(f"[debug] intermediate files: {tmp}")

    composers = {
        "cover": lambda i: compose_cover(
            project, vibe, report_type=args.report_type,
            period_start=period_start, period_end=period_end,
            funder_logo_uri=funder_logo_uri),
        "overview": lambda i: compose_overview(
            project, vibe, summary=summary,
            period_start=period_start, period_end=period_end),
        "impact": lambda i: compose_impact(project, vibe, idx=i),
        "deliverables": lambda i: compose_deliverables(project, vibe, idx=i),
        "activities": lambda i: compose_activities(
            project, vibe, idx=i, gallery_uris=gallery_uris),
        "stories": lambda i: compose_stories(
            project, vibe, idx=i, testimonials=testimonials),
        "geographic": lambda i: compose_geographic(project, vibe, idx=i),
        "financials": lambda i: compose_financials(project, vibe, idx=i),
        "way_forward": lambda i: compose_way_forward(
            project, vibe, idx=i, challenges=challenges, next_steps=next_steps),
        "closing": lambda i: compose_closing(project, vibe, idx=i),
    }

    try:
        # Compose every slide body in order, then stack into one HTML deck
        # consumed by dom-to-pptx via Puppeteer.
        slide_bodies: list[str] = []
        emitted: list[str] = []
        for idx, name in enumerate(manifest, start=1):
            composer = composers.get(name)
            if not composer:
                continue
            body_inner = composer(idx)
            if not body_inner:
                continue
            slide_bodies.append(body_inner)
            emitted.append(name)

        if not slide_bodies:
            raise RuntimeError("No slides emitted from manifest — data too sparse.")

        deck_path = tmp / "deck.html"
        deck_path.write_text(deck_html(theme, slide_bodies), encoding="utf-8")
        if keep:
            print(f"[debug] composed deck: {deck_path}")

        # Hand off to the Node export script. It uses puppeteer-core to drive
        # the system Chrome we already located, loads the deck, injects
        # dom-to-pptx, and writes a real editable .pptx.
        out = Path(args.output).resolve()
        node_script = SCRIPT_DIR / "dom_export.cjs"
        if not node_script.exists():
            raise RuntimeError(f"dom_export.cjs not found at {node_script}")
        # Resolve the project repo root that owns node_modules. The Node
        # script needs `puppeteer-core` and `dom-to-pptx`, both installed at
        # the repo root (the deer-flow sidecar runs out of the same checkout).
        repo_root = _find_repo_root_with_node_modules(SCRIPT_DIR)
        if repo_root is None:
            raise RuntimeError(
                "Could not locate a node_modules directory above the script. "
                "Run `npm install dom-to-pptx puppeteer-core` from the repo root first."
            )
        cmd = [
            _node_binary(),
            str(node_script),
            "--chrome", chrome,
            "--html", str(deck_path),
            "--out", str(out),
        ]
        try:
            proc = subprocess.run(
                cmd, capture_output=True, timeout=240, cwd=str(repo_root),
            )
        except subprocess.TimeoutExpired:
            raise RuntimeError("dom_export.cjs timed out after 240s")
        if proc.returncode != 0:
            stderr = proc.stderr.decode("utf-8", errors="ignore")
            stdout = proc.stdout.decode("utf-8", errors="ignore")
            raise RuntimeError(
                f"dom_export.cjs failed (exit {proc.returncode}).\n"
                f"--- stdout ---\n{stdout}\n--- stderr ---\n{stderr}"
            )

        if not out.exists():
            raise RuntimeError(f"Expected output file not written: {out}")
        size = out.stat().st_size
        if size < 8000:
            raise RuntimeError(f"Generated PPTX is suspiciously small ({size} bytes)")
        print(
            f"WROTE: {out} ({size:,} bytes, {len(emitted)} slides, "
            f"vibe={vibe.key}, theme={theme}, slides={'+'.join(emitted)})"
        )
        return 0
    finally:
        if not keep:
            try:
                shutil.rmtree(tmp, ignore_errors=True)
            except Exception:
                pass


if __name__ == "__main__":
    try:
        sys.exit(main())
    except Exception:
        traceback.print_exc()
        sys.exit(1)
